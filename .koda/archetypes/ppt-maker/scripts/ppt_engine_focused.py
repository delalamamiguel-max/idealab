#!/usr/bin/env python3
"""
PPT Engine - Focused Bulletproof Implementation
Building one slide type at a time with verification
"""

import sys
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass

try:
    import yaml
except ImportError:
    import json as yaml

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO


# ==============================================================================
# Constants
# ==============================================================================

# PowerPoint uses English Metric Units (EMUs) internally
# 1 inch = 914,400 EMUs
EMUS_PER_INCH = 914400

# Footer zone check (in inches) - validation still checks diagram doesn't extend into footer
FOOTER_ZONE_TOP_INCHES = 6.5


# ==============================================================================
# Configuration Data Structures
# ==============================================================================
@dataclass
class FieldConfig:
    """Configuration for a single field in a slide type"""
    char_limit: int
    description: Optional[str] = None

@dataclass
class SlideTypeConfig:
    """Configuration for a slide type"""
    name: str
    description: str
    fields: Dict[str, FieldConfig]
    diagram_support: bool = False
    required: bool = False
    position: Optional[str] = None
    best_for: Optional[List[str]] = None
    notes: Optional[str] = None

@dataclass
class ValidationRules:
    """Validation rules from configuration"""
    text_length_enabled: bool = True
    aspect_ratio_tolerance: float = 0.10
    diagram_scaling_threshold: float = 0.95
    diagram_coverage_threshold: float = 0.50
    diagram_safe_area_width: float = 9.5
    diagram_safe_area_height: float = 6.0


# ==============================================================================
# Slide Type Mappings - PowerPoint Implementation Details
# These contain PowerPoint-specific information (layout names, template indices)
# that are implementation details, not configuration data
# ==============================================================================

SLIDE_TYPE_MAPPINGS = {
    "01_title": {
        "layout_name": "Cover w/ Globe (White)",
        "template_slide_index": 0,  # Copy images from this template slide
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "metadata": "Text Placeholder 3"
        }
    },
    "02_cover_photo": {
        "layout_name": "Cover w/ Subtitle (White)",
        "template_slide_index": 1,  # Copy images from this template slide
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "presenter": "Text Placeholder 3",
            "date": "Text Placeholder 4"
        }
    },
    "03_divider": {
        "layout_name": "Divider 01",
        "template_slide_index": 2,  # Simple divider slide
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1"
        }
    },
    "04_divider_section": {
        "layout_name": "Divider 05 (w/ Slab #)",
        "template_slide_index": 3,  # Divider with section number and wave mask
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "section_number": "Text Placeholder 3"
        }
    },
    "05_content_single_column": {
        "layout_name": "Title + Subtitle (1 column w/ headers)",
        "template_slide_index": 4,  # Standard content slide with header and body
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "header": "Content Placeholder 3",
            "body": "Text Placeholder 4"
        }
    },
    "06_large_text_three_columns": {
        "layout_name": "Large Text Block w/ Blue Curve",
        "template_slide_index": 5,  # Large statement + 3 columns with wave mask
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "statement": "Content Placeholder 2",
            "column1_header": "Content Placeholder 3",
            "column1_body": "Text Placeholder 4",
            "column2_header": "Content Placeholder 5",
            "column2_body": "Text Placeholder 6",
            "column3_header": "Content Placeholder 7",
            "column3_body": "Text Placeholder 8"
        }
    },
    "07_two_columns": {
        "layout_name": "Title + Subtitle (2 column w/ headers)",
        "template_slide_index": 6,  # Title + subtitle with 2-column layout
        "shapes": {
            # Shape names when creating from LAYOUT
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "column1_header": "Content Placeholder 3",
            "column2_header": "Content Placeholder 4",
            "column1_body": "Text Placeholder 5",
            "column2_body": "Text Placeholder 6"
        }
    },
    "08_four_columns_enumerated": {
        "layout_name": "Title + Subtitle (4 column w/ headers)",
        "template_slide_index": 7,  # 4-column enumerated layout (numbers are in layout)
        "shapes": {
            # Shape names when creating from LAYOUT
            # Note: Numbers 1,2,3,4 are NOT placeholders - they're in the layout
            "title": "Title 1",
            "subtitle": "Text Placeholder 2",
            "column1_header": "Content Placeholder 3",
            "column2_header": "Content Placeholder 4",
            "column1_body": "Text Placeholder 5",
            "column2_body": "Text Placeholder 6",
            "column3_header": "Content Placeholder 7",
            "column3_body": "Text Placeholder 8",
            "column4_header": "Content Placeholder 9",
            "column4_body": "Text Placeholder 10"
        }
    },
    "09_grid_with_icons": {
        "layout_name": "Title + Subtitle (1/2 Blue Curve)",
        "template_slide_index": 8,  # Title + statement with 2x2 grid of icons/headers/bodies
        "shapes": {
            # Shape names when creating from LAYOUT
            # Note: Icons (Graphic 2,4,5,6) are NOT placeholders - they'll be copied
            "title": "Title 10",
            "statement": "Content Placeholder 9",
            "grid1_header": "Content Placeholder 1",
            "grid1_body": "Text Placeholder 2",
            "grid2_header": "Content Placeholder 3",
            "grid2_body": "Text Placeholder 4",
            "grid3_header": "Content Placeholder 5",
            "grid3_body": "Text Placeholder 6",
            "grid4_header": "Content Placeholder 7",
            "grid4_body": "Text Placeholder 8"
        }
    },
    "10_content_with_image": {
        "layout_name": "Title + Subtitle (1/2 Image on Right)",
        "template_slide_index": 9,  # Simple content on left with brand image on right
        "shapes": {
            # Shape names when creating from LAYOUT
            # Note: Picture placeholder (idx=19) will auto-fill from template
            "title": "Title 2",
            "subtitle": "Text Placeholder 3",
            "body": "Text Placeholder 4"
        }
    },
    "11_blank_end_slide": {
        "layout_name": "End Slide",
        "template_slide_index": 10,  # Blank slide with AT&T logo background
        "shapes": {
            # No editable content - logo and background are in the layout
        }
    },
    "12_final_slide": {
        "layout_name": "Logo End Slide (White)",
        "template_slide_index": 11,  # Final blank slide with centered AT&T logo
        "shapes": {
            # No editable content - centered logo is in the layout
            # This should always be the last slide
        }
    }
}


class PPTEngine:
    def __init__(self, template_path: str, assets_folder: Optional[str] = None, config_path: Optional[str] = None):
        self.template_path = Path(template_path)
        
        # If template path is relative and doesn't exist in CWD, try resolving from script location
        if not self.template_path.is_absolute() and not self.template_path.exists():
            # Get the script's directory (.cdo-aifc/scripts/python/)
            script_dir = Path(__file__).resolve().parent
            
            # Remove .cdo-aifc/ prefix if present in template_path (for default path compatibility)
            template_path_clean = str(template_path)
            if template_path_clean.startswith('.cdo-aifc/'):
                template_path_clean = template_path_clean[len('.cdo-aifc/'):]
            
            # Navigate up to .cdo-aifc/ and append the cleaned path
            potential_path = script_dir.parent.parent / template_path_clean
            if potential_path.exists():
                self.template_path = potential_path
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}\nSearched in CWD and script-relative locations")
        
        # Assets folder for persisting YAML and images
        self.assets_folder = Path(assets_folder) if assets_folder else Path("assets")
        self.assets_folder.mkdir(exist_ok=True)
        
        # Load configuration from YAML
        self._load_configuration(config_path)
        
        # Load template to get layouts and template slides
        self.template = Presentation(str(self.template_path))
        
        # Extract theme colors from configuration
        self._load_theme_colors()
    
    def _load_configuration(self, config_path: Optional[str] = None):
        """Load configuration from YAML file and parse into structured objects"""
        if config_path is None:
            # Default: look for config file in script directory
            script_dir = Path(__file__).resolve().parent
            config_path = script_dir / "ppt_engine_focused.yaml"
        else:
            config_path = Path(config_path)
        
        if not config_path.exists():
            raise FileNotFoundError(f"Configuration file not found: {config_path}")
        
        with open(config_path, 'r') as f:
            raw_config = yaml.safe_load(f)
        
        # Parse slide types into structured objects
        self.slide_types = {}
        for type_id, type_data in raw_config.get('slide_types', {}).items():
            # Parse fields
            fields = {}
            for field_name, field_data in type_data.get('fields', {}).items():
                fields[field_name] = FieldConfig(
                    char_limit=field_data.get('char_limit', 9999),
                    description=field_data.get('description')
                )
            
            # Create slide type config
            self.slide_types[type_id] = SlideTypeConfig(
                name=type_data.get('name', type_id),
                description=type_data.get('description', ''),
                fields=fields,
                diagram_support=type_data.get('diagram_support', False),
                required=type_data.get('required', False),
                position=type_data.get('position'),
                best_for=type_data.get('best_for'),
                notes=type_data.get('notes')
            )
        
        # Parse validation rules
        val_rules = raw_config.get('validation_rules', {})
        diagram_sizing = val_rules.get('diagram_sizing', {})
        safe_area = diagram_sizing.get('safe_area', {})
        
        self.validation = ValidationRules(
            text_length_enabled=val_rules.get('text_length', {}).get('enabled', True),
            aspect_ratio_tolerance=val_rules.get('aspect_ratio', {}).get('tolerance', 0.10),
            diagram_scaling_threshold=diagram_sizing.get('scaling_threshold', 0.95),
            diagram_coverage_threshold=diagram_sizing.get('coverage_threshold', 0.50),
            diagram_safe_area_width=safe_area.get('width_inches', 9.5),
            diagram_safe_area_height=safe_area.get('height_inches', 6.0)
        )
        
        # Keep theme colors as raw dict for now (already used elsewhere)
        self.theme_colors_config = raw_config.get('theme_colors', {})
    
    def _load_theme_colors(self):
        """Extract theme colors from configuration"""
        from pptx.dml.color import RGBColor
        
        # Load colors from config
        self.theme_colors = {}
        for color_name, color_config in self.theme_colors_config.items():
            rgb = color_config.get('rgb', [0, 0, 0])
            self.theme_colors[color_name] = RGBColor(rgb[0], rgb[1], rgb[2])
    
    def validate_image_aspect_ratio(self, image_path: str, slide_width_inches: float, 
                                    slide_height_inches: float, tolerance: float = 0.1) -> bool:
        """
        Validate that image won't be distorted when displayed on slide.
        
        Compares SOURCE aspect ratio (image file) vs DISPLAY aspect ratio (on slide).
        If they differ by more than tolerance, the image will appear stretched/squashed.
        
        Args:
            image_path: Path to source image file
            slide_width_inches: Width on slide (inches)
            slide_height_inches: Height on slide (inches)
            tolerance: Acceptable distortion (0.1 = 10%)
        
        Returns:
            True if distortion is acceptable
        
        Raises:
            ValueError: If image would be distorted > tolerance
        """
        from PIL import Image
        
        # Load image to get source dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
            source_aspect = img_width / img_height
        
        # Calculate display aspect ratio
        display_aspect = slide_width_inches / slide_height_inches
        
        # Calculate distortion percentage
        distortion = abs(source_aspect - display_aspect) / source_aspect
        
        if distortion <= tolerance:
            print(f"  ✓ Image aspect ratio OK: {distortion*100:.1f}% distortion (acceptable)")
            return True
        else:
            # Calculate recommended dimensions that preserve source aspect
            option1_height = slide_width_inches / source_aspect
            option2_width = slide_height_inches * source_aspect
            
            error_msg = (
                f"❌ Image will be distorted by {distortion*100:.1f}%\n"
                f"   Source: {img_width}×{img_height} (aspect: {source_aspect:.2f}:1)\n"
                f"   Display: {slide_width_inches:.2f}\"×{slide_height_inches:.2f}\" "
                f"(aspect: {display_aspect:.2f}:1)\n"
                f"   Recommended dimensions to preserve aspect ratio:\n"
                f"     Option 1: {slide_width_inches:.2f}\"×{option1_height:.2f}\" (fit to width)\n"
                f"     Option 2: {option2_width:.2f}\"×{slide_height_inches:.2f}\" (fit to height)\n"
                f"   Alternatively: Set validate_aspect_ratio: false to accept distortion"
            )
            raise ValueError(error_msg)
    
    def validate_and_adjust_dimensions(self, image_path: str, width: float, height: float, 
                                      tolerance: float = 0.10) -> tuple:
        """
        Validate aspect ratio and auto-adjust dimensions if needed to prevent distortion.
        
        Args:
            image_path: Path to source image file
            width: Requested width on slide (inches)
            height: Requested height on slide (inches)
            tolerance: Acceptable distortion (0.10 = 10%)
        
        Returns:
            (adjusted_width, adjusted_height, was_adjusted) tuple
        """
        from PIL import Image
        
        # Load image to get source dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size
            source_aspect = img_width / img_height
        
        display_aspect = width / height
        distortion = abs(source_aspect - display_aspect) / source_aspect
        
        if distortion <= tolerance:
            print(f"  ✓ Aspect ratio OK: {distortion*100:.1f}% distortion (acceptable)")
            return (width, height, False)
        
        # Auto-adjust: preserve width, calculate height from source aspect
        new_height = width / source_aspect
        new_width = width
        
        # Check if fits in safe zone (avoid footer/logo at bottom)
        if new_height > 6.0:
            # Too tall, fit to height instead
            new_height = 6.0
            new_width = new_height * source_aspect
            print(f"  ⚠️  Auto-adjusted to fit safe zone: {new_width:.2f}\"×{new_height:.2f}\"")
        else:
            print(f"  ⚠️  Auto-adjusted dimensions to preserve aspect ratio:")
        
        print(f"     Original request: {width:.2f}\"×{height:.2f}\" (would distort {distortion*100:.1f}%)")
        print(f"     Adjusted result: {new_width:.2f}\"×{new_height:.2f}\" (preserves {source_aspect:.2f}:1)")
        
        return (new_width, new_height, True)
    
    def detect_shape_collisions(self, slide) -> List[Dict[str, Any]]:
        """
        Detect overlapping or colliding shapes on a slide.
        
        Args:
            slide: The slide to check for collisions
            
        Returns:
            List of collision dictionaries with shape details
        """
        collisions = []
        shapes = list(slide.shapes)
        
        for i, shape1 in enumerate(shapes):
            # Skip if shape doesn't have position/size (like group shapes)
            if not hasattr(shape1, 'left') or not hasattr(shape1, 'top'):
                continue
                
            for shape2 in shapes[i+1:]:
                if not hasattr(shape2, 'left') or not hasattr(shape2, 'top'):
                    continue
                
                # Check if bounding boxes overlap
                if self._shapes_overlap(shape1, shape2):
                    collision = {
                        'shape1': shape1.name,
                        'shape2': shape2.name,
                        'shape1_bounds': {
                            'left': shape1.left,
                            'top': shape1.top,
                            'width': shape1.width,
                            'height': shape1.height
                        },
                        'shape2_bounds': {
                            'left': shape2.left,
                            'top': shape2.top,
                            'width': shape2.width,
                            'height': shape2.height
                        }
                    }
                    collisions.append(collision)
        
        return collisions
    
    def _shapes_overlap(self, shape1, shape2, margin_emu: int = 0) -> bool:
        """
        Check if two shapes' bounding boxes overlap.
        
        Args:
            shape1: First shape
            shape2: Second shape
            margin_emu: Optional margin in EMUs (English Metric Units) to add as buffer
            
        Returns:
            True if shapes overlap, False otherwise
        """
        # Get bounding boxes with optional margin
        s1_left = shape1.left - margin_emu
        s1_right = shape1.left + shape1.width + margin_emu
        s1_top = shape1.top - margin_emu
        s1_bottom = shape1.top + shape1.height + margin_emu
        
        s2_left = shape2.left - margin_emu
        s2_right = shape2.left + shape2.width + margin_emu
        s2_top = shape2.top - margin_emu
        s2_bottom = shape2.top + shape2.height + margin_emu
        
        # Check for overlap (NOT touching is sufficient separation)
        horizontal_overlap = s1_left < s2_right and s1_right > s2_left
        vertical_overlap = s1_top < s2_bottom and s1_bottom > s2_top
        
        return horizontal_overlap and vertical_overlap
    
    def add_diagram_shapes(self, slide, diagram_type: str, config: Dict[str, Any]):
        """
        Add programmatic diagram shapes to a slide with proper colors and sizing.
        
        Args:
            slide: The slide to add shapes to
            diagram_type: Type of diagram ('flow', 'timeline', 'custom')
            config: Configuration dictionary with diagram parameters
        """
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        
        if diagram_type == 'flow':
            self._add_flow_diagram(slide, config)
        elif diagram_type == 'timeline':
            self._add_timeline_diagram(slide, config)
        else:
            print(f"  ⚠️  Unknown diagram type: {diagram_type}")
    
    def _add_flow_diagram(self, slide, config: Dict[str, Any]):
        """Add a flow diagram with proper sizing and colors"""
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        
        steps = config.get('steps', [])
        if not steps:
            return
        
        # Calculate appropriate box size based on text length
        max_text_length = max(len(step.get('text', '')) for step in steps)
        
        # Adjust box width based on text length (with reasonable min/max)
        box_width = max(1.0, min(2.0, max_text_length / 10))
        box_height = 0.7
        
        # Use theme colors
        fill_color = self.theme_colors['primary_blue']
        line_color = self.theme_colors['dark_blue']
        text_color = self.theme_colors['text_light']
        
        for step in steps:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(step.get('x', 1.0)),
                Inches(step.get('y', 2.0)),
                Inches(box_width),
                Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.color.rgb = line_color
            shape.line.width = Pt(2.5)
            
            # Add text with appropriate sizing
            text_frame = shape.text_frame
            text = step.get('text', '')
            text_frame.text = text
            
            # Calculate font size based on text length to prevent wrapping
            text_length = len(text)
            if text_length < 15:
                font_size = 12
            elif text_length < 25:
                font_size = 11
            else:
                font_size = 10
            
            text_frame.paragraphs[0].font.size = Pt(font_size)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = text_color
            text_frame.paragraphs[0].alignment = 1  # Center
            text_frame.word_wrap = True
        
        # Add connectors if specified
        connectors = config.get('connectors', [])
        for conn in connectors:
            connector = slide.shapes.add_connector(
                1,  # Straight line
                Inches(conn.get('x1', 0)),
                Inches(conn.get('y1', 0)),
                Inches(conn.get('x2', 1)),
                Inches(conn.get('y2', 1))
            )
            connector.line.color.rgb = line_color
            connector.line.width = Pt(2)
        
        print(f"  ✓ Added flow diagram with {len(steps)} steps")
    
    def _add_timeline_diagram(self, slide, config: Dict[str, Any]):
        """Add a timeline diagram with proper sizing and colors"""
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        
        phases = config.get('phases', [])
        if not phases:
            return
        
        # Use theme colors
        fill_color = self.theme_colors['light_blue']
        line_color = self.theme_colors['dark_blue']
        text_color = self.theme_colors['dark_blue']
        
        # Main timeline bar
        timeline_config = config.get('timeline_bar', {})
        timeline = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(timeline_config.get('x', 1.0)),
            Inches(timeline_config.get('y', 3.3)),
            Inches(timeline_config.get('width', 9.0)),
            Inches(timeline_config.get('height', 0.15))
        )
        timeline.fill.solid()
        timeline.fill.fore_color.rgb = line_color
        timeline.line.width = Pt(0)
        
        # Add phase boxes and milestone circles
        for phase in phases:
            # Calculate box size based on text length
            text_length = len(phase.get('name', ''))
            box_width = max(1.2, min(2.0, text_length / 10))  # Adjusted for better spacing
            phase_x = phase.get('x', 1.0) - (box_width / 2)  # Center box on x position
            phase_y = phase.get('y', 2.0)
            
            # Add phase box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(phase_x),
                Inches(phase_y),
                Inches(box_width),
                Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = fill_color
            box.line.color.rgb = line_color
            box.line.width = Pt(2)
            
            # Add text
            text_frame = box.text_frame
            text_frame.text = phase.get('name', '')
            p = text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = text_color
            p.alignment = 1
            
            # Add milestone circle on timeline (centered at phase position)
            circle_x = phase.get('x', 1.0) - 0.15  # Center at original x position
            circle_y = timeline_config.get('y', 3.3) - 0.075  # Center on timeline
            
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(circle_x),
                Inches(circle_y),
                Inches(0.3),
                Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme_colors['primary_blue']
            circle.line.color.rgb = line_color
            circle.line.width = Pt(3)
        
        print(f"  ✓ Added timeline diagram with {len(phases)} phases and milestone markers")
    
    def _process_diagram(self, slide, diagram_spec: Dict[str, Any]):
        """
        Process diagram specification and add to slide.
        
        Supports two types:
        1. Image-based: diagram from external file (e.g., Mermaid CLI output)
        2. Programmatic: diagram generated with python-pptx shapes
        
        Args:
            slide: The slide to add diagram to
            diagram_spec: Diagram specification dictionary
        """
        from pptx.util import Inches
        
        diagram_type = diagram_spec.get('type')
        
        if diagram_type == 'image':
            # Image-based diagram (e.g., from Mermaid CLI)
            image_path = diagram_spec.get('path')
            if not image_path:
                print(f"  ⚠️  Image diagram missing 'path' field")
                return
            
            # Resolve path relative to assets folder if not absolute
            if not Path(image_path).is_absolute():
                image_path = str(self.assets_folder / image_path)
            
            if not Path(image_path).exists():
                print(f"  ⚠️  Image not found: {image_path}")
                return
            
            # Get position and size with smart defaults
            position = diagram_spec.get('position', {})
            
            # Check if this is a blank/diagram-only slide
            slide_layout = slide.slide_layout.name
            is_diagram_slide = any(name in slide_layout for name in ["Blank", "End Slide"])
            
            # Auto-size logic: preserve source aspect ratio if no position specified
            auto_sized_blank_slide = False  # Track if we auto-sized for blank slide
            if not position and is_diagram_slide:
                try:
                    from PIL import Image
                    with Image.open(image_path) as img:
                        img_width, img_height = img.size
                        source_aspect = img_width / img_height
                    
                    # For blank/end slides: center and maximize space usage
                    if "Blank" in slide_layout or "End Slide" in slide_layout:
                        auto_sized_blank_slide = True  # Mark as auto-sized
                        
                        # Use configured safe area from validation config
                        max_width = self.validation.diagram_safe_area_width
                        max_height = self.validation.diagram_safe_area_height
                        
                        # Calculate dimensions preserving aspect ratio
                        if source_aspect > (max_width / max_height):
                            # Wide diagram - fit to width
                            width = max_width
                            height = width / source_aspect
                        else:
                            # Tall diagram - fit to height
                            height = max_height
                            width = height * source_aspect
                        
                        # Position: horizontally centered, fixed top margin
                        x = (13.33 - width) / 2  # Center horizontally on 13.33" wide slide
                        y = 0.5  # Fixed 0.5" top margin (leaves 1.5" bottom for footer)
                        
                        print(f"  ℹ️  Centered diagram on blank slide: {width:.2f}\"×{height:.2f}\" at ({x:.2f}\", {y:.2f}\")")
                        
                        # Analyze diagram fit: check area coverage first, then scaling ratio
                        actual_area = width * height
                        max_possible_area = max_width * max_height
                        area_coverage = actual_area / max_possible_area
                        scaling_ratio = min(width / max_width, height / max_height)
                        
                        # Check 1: Area coverage - warn if using less than threshold
                        if area_coverage < self.validation.diagram_coverage_threshold:
                            print(f"  ⚠️  Low area coverage: Diagram uses only {area_coverage*100:.0f}% of available slide space")
                            print(f"     This may indicate an inefficient aspect ratio for the slide dimensions")
                            print(f"     Consider adjusting diagram proportions or content layout")
                        
                        # Check 2: Scaling ratio - warn if below threshold
                        if scaling_ratio < self.validation.diagram_scaling_threshold:
                            print(f"  ⚠️  Small scaling: Diagram scaled to {scaling_ratio*100:.0f}% of available space")
                            print(f"     Text and details may be difficult to read at this size")
                            print(f"     Consider regenerating with larger dimensions for better visibility")
                            print(f"     Recommended source dimensions: {int(max_width * 250)}×{int(max_height * 250)} pixels")
                            # For mermaid diagrams, could offer automatic regeneration
                            if '.mmd' in str(image_path) or 'mermaid' in str(image_path).lower():
                                print(f"     💡 This appears to be a Mermaid diagram - consider regenerating with:")
                                print(f"        mmdc -i {Path(image_path).stem}.mmd -o {Path(image_path).stem}.png -w {int(max_width * 250)} -H {int(max_height * 250)}")
                                print(f"        # Note: Adjusted for safe zone (footer at 6.5\"), max diagram height is 6.0\"")
                    else:
                        # Standard content slide auto-sizing
                        x = 0.5
                        width = 9.0
                        
                        # Calculate height to preserve aspect ratio
                        height = width / source_aspect
                        
                        # Check if fits in safe zone (leave room for footer/logo)
                        if height <= 5.75:
                            y = 0.75
                            print(f"  ℹ️  Auto-sized to preserve aspect ratio: {width}\"×{height:.2f}\" (source: {source_aspect:.2f}:1)")
                        else:
                            # Too tall, fit to height instead
                            y = 0.5
                            height = 5.75
                            width = height * source_aspect
                            print(f"  ℹ️  Auto-sized to fit safe zone: {width:.2f}\"×{height}\" (source: {source_aspect:.2f}:1)")
                    
                except Exception as e:
                    # Fall back to defaults if auto-sizing fails
                    print(f"  ⚠️  Could not auto-size image: {e}")
                    x = 0.5
                    y = 0.75
                    width = 9.0
                    height = 5.5
            else:
                # Manual positioning specified (NOT RECOMMENDED)
                x = position.get('x', 0.5)
                y = position.get('y', 0.75)
                width = position.get('width', 9.0)
                height = position.get('height', 5.5)
                
                print(f"  ⚠️  Manual positioning detected: {x:.2f}\", {y:.2f}\", {width:.2f}\"×{height:.2f}\"")
                print(f"     Manual positions may cause footer overlap and skip readability checks")
                
                # Initial footer overlap check (will recheck after aspect ratio adjustment)
                if y + height > 6.5:
                    print(f"  ❌  Manual position causes footer overlap: y({y:.2f}) + height({height:.2f}) = {y+height:.2f} > 6.5\"")
                    print(f"     Consider omitting 'position' to let engine auto-size safely")
            
            # Validate and auto-adjust aspect ratio (enabled by default)
            # Skip validation if we just auto-sized for blank slide (already perfect aspect ratio)
            # User must explicitly set validate_aspect_ratio: false to disable
            if diagram_spec.get('validate_aspect_ratio', True) and not auto_sized_blank_slide:
                width, height, was_adjusted = self.validate_and_adjust_dimensions(
                    image_path, width, height, tolerance=0.10
                )
                
                # CRITICAL: Recheck footer overlap after aspect ratio adjustment
                if was_adjusted and y + height > 6.5:
                    print(f"  ❌  FOOTER OVERLAP after aspect ratio adjustment: y({y:.2f}) + height({height:.2f}) = {y+height:.2f} > 6.5\"")
                    # Adjust to fit within safe zone
                    max_safe_height = 6.5 - y
                    if max_safe_height > 1.0:  # Only adjust if reasonable space available
                        height = max_safe_height
                        print(f"     Auto-adjusted height to {height:.2f}\" to avoid footer")
                    else:
                        # Position is too low, adjust y as well
                        y = 0.5
                        height = 6.0
                        print(f"     Repositioned to y={y:.2f}\", height={height:.2f}\" to fit safely")
                        
            elif auto_sized_blank_slide:
                print(f"  ✓  Skipping aspect ratio validation (already preserved by auto-sizing)")
            
            # Check scaling ratio for manual positions (for readability)
            if position:
                try:
                    from PIL import Image
                    with Image.open(image_path) as img:
                        img_width, img_height = img.size
                        source_aspect = img_width / img_height
                    
                    # For blank slides, check against max safe area
                    if is_diagram_slide:
                        max_width = self.validation.diagram_safe_area_width
                        max_height = self.validation.diagram_safe_area_height
                        scaling_ratio = min(width / max_width, height / max_height)
                        
                        if scaling_ratio < self.validation.diagram_scaling_threshold:
                            print(f"  ⚠️  Manual position results in small diagram: {scaling_ratio*100:.0f}% of available space")
                            print(f"     Text may be difficult to read at this size")
                            print(f"     Recommendation: Remove 'position' block to enable auto-sizing")
                except Exception as e:
                    print(f"  ⚠️  Could not check scaling ratio: {e}")
            
            # Add image to slide with error handling
            try:
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches(x),
                    Inches(y),
                    Inches(width),
                    Inches(height)
                )
                
                # Verify image was actually added
                if pic.width == 0 or pic.height == 0:
                    raise ValueError(f"Image embedded but has zero dimensions")
                
                print(f"  ✓ Embedded image: {Path(image_path).name} ({width:.2f}\"×{height:.2f}\")")
                
            except Exception as e:
                error_msg = f"Failed to embed image {image_path}: {str(e)}"
                print(f"  ❌ {error_msg}")
                raise ValueError(error_msg)
            
        elif diagram_type == 'programmatic':
            # Programmatic diagram (python-pptx shapes)
            style = diagram_spec.get('style')
            config = diagram_spec.get('config', {})
            
            if style == 'flow':
                self._add_flow_diagram(slide, config)
            elif style == 'timeline':
                self._add_timeline_diagram(slide, config)
            else:
                print(f"  ⚠️  Unknown programmatic diagram style: {style}")
        
        else:
            print(f"  ⚠️  Unknown diagram type: {diagram_type}")
    
    def create_from_yaml_spec(self, yaml_spec: str, output_path: str, persist_yaml: bool = True) -> str:
        """Create presentation from YAML specification"""
        spec = yaml.safe_load(yaml_spec)
        
        # Optionally persist YAML to assets folder for reuse
        if persist_yaml:
            output_name = Path(output_path).stem
            yaml_path = self.assets_folder / f"{output_name}.yaml"
            with open(yaml_path, 'w') as f:
                f.write(yaml_spec)
            print(f"  📄 Saved YAML spec to: {yaml_path}")
        
        return self._create_from_spec(spec, output_path)
    
    def _validate_constitutional_rules(self, spec: Dict):
        """
        Validate constitutional hard-stop rules before generating presentation.
        
        Enforces:
        - Required slides: 01_title (first), 12_final_slide (last)
        - No placeholder text in content
        
        Raises:
            ValueError: If constitutional rules are violated
        """
        errors = []
        
        # Check for mandatory title slide
        if 'title' not in spec:
            errors.append("❌ CONSTITUTIONAL VIOLATION: Missing mandatory title slide (type '01_title')")
        
        # Check for mandatory final slide
        slides = spec.get('slides', [])
        if not slides:
            errors.append("❌ CONSTITUTIONAL VIOLATION: No slides defined in spec")
        else:
            last_slide_type = slides[-1].get('type')
            if last_slide_type != '12_final_slide':
                errors.append(
                    f"❌ CONSTITUTIONAL VIOLATION: Last slide must be type '12_final_slide', "
                    f"found '{last_slide_type}'"
                )
        
        # Check for placeholder text (common patterns)
        placeholder_patterns = [
            'lorem ipsum',
            '[insert',
            'placeholder',
            'todo',
            'tbd',
            'xxx'
        ]
        
        def check_content(content_dict, slide_ref):
            """Recursively check content for placeholder text"""
            if isinstance(content_dict, dict):
                for key, value in content_dict.items():
                    if isinstance(value, str):
                        value_lower = value.lower()
                        for pattern in placeholder_patterns:
                            if pattern in value_lower:
                                errors.append(
                                    f"⚠️  WARNING: Possible placeholder text in {slide_ref}.{key}: '{value[:50]}...'"
                                )
                    elif isinstance(value, dict):
                        check_content(value, f"{slide_ref}.{key}")
        
        # Check title slide content
        if 'title' in spec:
            check_content(spec['title'], 'title')
        
        # Check all slides
        for i, slide_spec in enumerate(slides):
            content = slide_spec.get('content', {})
            check_content(content, f'slides[{i}]')
        
        # Raise error if critical violations found
        critical_errors = [e for e in errors if 'CONSTITUTIONAL VIOLATION' in e]
        if critical_errors:
            error_msg = "\n".join(critical_errors)
            raise ValueError(f"\n\n{error_msg}\n\nPresentation generation aborted due to constitutional violations.")
        
        # Print warnings for non-critical issues
        warnings = [e for e in errors if 'WARNING' in e]
        for warning in warnings:
            print(f"  {warning}")
    
    def _validate_content_lengths(self, spec: Dict[str, Any]) -> List[str]:
        """
        Validate content lengths BEFORE rendering.
        Prevents silent truncation and ellipses.
        
        Returns list of validation errors.
        """
        errors = []
        
        for idx, slide_spec in enumerate(spec.get('slides', []), start=2):
            slide_type = slide_spec.get('type')
            content = slide_spec.get('content', {})
            
            # Get constraints for this slide type from dataclass
            slide_config = self.slide_types.get(slide_type)
            if not slide_config:
                continue
            
            constraints = {name: field.char_limit for name, field in slide_config.fields.items()}
            
            for field_name, field_value in content.items():
                if not isinstance(field_value, str):
                    continue
                
                max_length = constraints.get(field_name)
                if max_length and len(field_value) > max_length:
                    errors.append(
                        f"Slide {idx} ({slide_type}): Field '{field_name}' exceeds {max_length} char limit "
                        f"({len(field_value)} chars). Text must be rewritten to fit.\n"
                        f"   Content: '{field_value[:100]}...'"
                    )
        
        return errors
    
    def _create_from_spec(self, spec: Dict, output_path: str) -> str:
        """Create new presentation from spec"""
        # Content length validation: Fail fast if text exceeds limits
        length_errors = self._validate_content_lengths(spec)
        if length_errors:
            error_msg = "❌ Content length validation FAILED:\n\n" + "\n\n".join(length_errors)
            error_msg += "\n\n💡 Text must be rewritten to fit within character limits. No ellipses truncation allowed."
            print(error_msg)
            raise ValueError(error_msg)
        
        # Constitutional validation: Enforce mandatory slides
        self._validate_constitutional_rules(spec)
        
        # Use the template presentation object that was already loaded
        prs = self.template
        
        # Track original slide count
        original_slide_count = len(prs.slides)
        
        # Extract footer information from spec
        footer_info = self._extract_footer_info(spec)
        
        # Add title slide if specified
        if 'title' in spec:
            self._add_title_slide(prs, spec['title'], footer_info)
        
        # Add content slides
        for slide_spec in spec.get('slides', []):
            slide_type = slide_spec.get('type')
            
            # Skip 01_title if already added via top-level 'title' key (prevents duplicates)
            if slide_type == '01_title' and 'title' in spec:
                print(f"⚠️  Skipping duplicate title slide in slides array (already added from top-level 'title' key)")
                continue
            
            content = slide_spec.get('content', {})
            diagram = slide_spec.get('diagram', None)  # Optional diagram specification
            
            # Use generic method for all slide types
            self._add_slide(prs, slide_type, content, footer_info, diagram)
        
        # Now remove the original template slides (from the beginning)
        # Our new slides are at the end, template slides are at the beginning
        for _ in range(original_slide_count):
            # Remove first slide (original template slides)
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Save
        prs.save(output_path)
        return output_path
    
    def _extract_footer_info(self, spec: Dict) -> Dict[str, str]:
        """Extract footer information from spec or use defaults"""
        from datetime import datetime
        
        # Get current date for defaults
        now = datetime.now()
        
        # Extract from spec or use defaults
        footer_info = {
            'title': spec.get('footer', {}).get('title', spec.get('title', {}).get('title', 'Presentation Title')),
            'month': spec.get('footer', {}).get('month', now.strftime('%B')),
            'day': spec.get('footer', {}).get('day', now.strftime('%d')),
            'year': spec.get('footer', {}).get('year', str(now.year))
        }
        
        return footer_info
    
    def _validate_slide_content_match(self, slide_type: str, content: Dict[str, str]) -> List[str]:
        """
        Validate that content matches the chosen template.
        Warns about potential mismatches.
        
        Returns:
            List of warning messages
        """
        warnings = []
        
        if slide_type == "09_grid_with_icons":
            # Check if grid content is actually provided
            has_grid_content = any([
                content.get(f"grid{i}_header") or content.get(f"grid{i}_body")
                for i in range(1, 5)
            ])
            
            if not has_grid_content:
                warnings.append(
                    "⚠️  Template '09_grid_with_icons' chosen but no grid content provided.\n"
                    "   This template has decorative icons (?, ▲, ✓, 🔔) that will appear even if empty.\n"
                    "   Consider using '06_large_text_three_columns' or '07_two_columns' instead.\n"
                    "   If icons are intentional, this is not an error."
                )
        
        if slide_type == "08_four_columns_enumerated":
            # Check if 4 columns are actually provided
            column_content = [
                content.get(f"column{i}_header") or content.get(f"column{i}_body")
                for i in range(1, 5)
            ]
            columns_provided = sum(1 for c in column_content if c)
            
            if columns_provided < 4:
                warnings.append(
                    f"⚠️  Template '08_four_columns_enumerated' expects 4 columns but {columns_provided} provided.\n"
                    "   Consider using '06_large_text_three_columns' (3 cols) or '07_two_columns' (2 cols)."
                )
        
        return warnings
    
    def _add_title_slide(self, prs: Presentation, content: Dict[str, str], footer_info: Dict[str, str]):
        """Add title slide with content - delegates to generic method"""
        self._add_slide(prs, "01_title", content, footer_info, diagram=None)
    
    def _add_slide(self, prs: Presentation, slide_type: str, content: Dict[str, str], 
                   footer_info: Dict[str, str], diagram: Optional[Dict[str, Any]] = None):
        """Generic method to add any slide type - creates from layout, copies images, updates text, adds diagrams"""
        if slide_type not in SLIDE_TYPE_MAPPINGS:
            print(f"  ⚠️  Unknown slide type: {slide_type}")
            return
        
        config = SLIDE_TYPE_MAPPINGS[slide_type]
        template_slide_idx = config.get("template_slide_index")
        
        if template_slide_idx is None:
            print(f"  ⚠️  Template slide index not found for: {slide_type}")
            return
        
        # Validate template/content match and warn if mismatch detected
        warnings = self._validate_slide_content_match(slide_type, content)
        for warning in warnings:
            print(f"  {warning}")
        
        # Get the template slide from THE SAME presentation object
        # (not from self.template which is a different object)
        template_slide = prs.slides[template_slide_idx]
        
        # Create new slide with same layout
        slide_layout = template_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Extract and copy images from template slide IN THE SAME PRESENTATION
        self._copy_images_from_template(template_slide, new_slide)
        
        # Copy footer and slide number placeholders from template (with updated text)
        self._copy_footer_elements(template_slide, new_slide, footer_info)
        
        # Copy non-placeholder text boxes (like enumeration numbers)
        self._copy_non_placeholder_text_boxes(template_slide, new_slide)
        
        # Copy placeholder positions from template (they may be adjusted from layout defaults)
        self._copy_placeholder_positions(template_slide, new_slide)
        
        # Update text content with formatting preservation
        self._update_text_with_formatting(template_slide, new_slide, slide_type, config["shapes"], content)
        
        # Add diagram if specified
        if diagram:
            self._process_diagram(new_slide, diagram)
    
    def _update_text_with_formatting(self, template_slide, new_slide, slide_type: str, shape_mapping: Dict[str, str], content: Dict[str, str]):
        """Update text content while preserving font formatting from template"""
        # Get length constraints for this slide type from dataclass
        slide_config = self.slide_types.get(slide_type)
        if not slide_config:
            return  # No config for this type
        
        constraints = {name: field.char_limit for name, field in slide_config.fields.items()}
        
        # Build maps by placeholder index (this matches template to new slide)
        template_shapes_by_idx = {}
        for shape in template_slide.shapes:
            if hasattr(shape, 'text_frame'):
                try:
                    idx = shape.placeholder_format.idx
                    template_shapes_by_idx[idx] = shape
                except ValueError:
                    # Not a placeholder, skip
                    continue
        
        new_shapes_by_idx = {}
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                try:
                    idx = shape.placeholder_format.idx
                    new_shapes_by_idx[idx] = shape
                except ValueError:
                    # Not a placeholder, skip
                    continue
        
        for content_key, shape_name in shape_mapping.items():
            if content_key not in content or not content[content_key]:
                continue
            
            # Check length constraint for this field
            text_content = content[content_key]
            if content_key in constraints:
                max_length = constraints[content_key]
                if len(text_content) > max_length:
                    original_length = len(text_content)
                    text_content = text_content[:max_length-3] + "..."
                    print(f"  ⚠️  Truncated '{content_key}' from {original_length} to {max_length} chars: '{text_content}'")
            
            # Find the shape in new slide by name
            target_shape = None
            target_idx = None
            for shape in new_slide.shapes:
                if shape.name == shape_name and shape.has_text_frame:
                    target_shape = shape
                    try:
                        target_idx = shape.placeholder_format.idx
                    except ValueError:
                        # Not a placeholder, that's okay
                        pass
                    break
            
            if not target_shape:
                continue
            
            # Get font properties from template shape with same placeholder index
            template_font = None
            if target_idx and target_idx in template_shapes_by_idx:
                template_shape = template_shapes_by_idx[target_idx]
                if len(template_shape.text_frame.paragraphs) > 0:
                    p = template_shape.text_frame.paragraphs[0]
                    if len(p.runs) > 0:
                        template_font = p.runs[0].font
            
            # Set the text with proper formatting
            text_frame = target_shape.text_frame
            text_frame.clear()  # Clear completely
            p = text_frame.paragraphs[0]
            p.text = text_content
            
            # Apply template font properties to the run
            if len(p.runs) > 0 and template_font:
                run = p.runs[0]
                if template_font.name:
                    run.font.name = template_font.name
                if template_font.size:
                    run.font.size = template_font.size
                if template_font.bold is not None:
                    run.font.bold = template_font.bold
                if template_font.color and template_font.color.type:
                    run.font.color.rgb = template_font.color.rgb
            
            print(f"  ✓ Set {shape_name} = {text_content[:30]}...")
    
    def _copy_images_from_template(self, template_slide, new_slide):
        """Extract images from template slide and fill picture placeholders in new slide"""
        # Find all images in template
        template_images = []
        for shape in template_slide.shapes:
            if hasattr(shape, 'image'):
                try:
                    template_images.append({
                        'name': shape.name,
                        'image': shape.image,
                        'image_bytes': shape.image.blob
                    })
                except:
                    continue
        
        if not template_images:
            return
        
        # Find picture placeholders in new slide and fill them
        for shape in new_slide.placeholders:
            if hasattr(shape, 'placeholder_format'):
                ph_type = str(shape.placeholder_format.type)
                if 'PICTURE' in ph_type and template_images:
                    try:
                        # Get the first available image
                        img_data = template_images.pop(0)
                        image_stream = BytesIO(img_data['image_bytes'])
                        
                        # Fill the placeholder with the image
                        # This preserves the placeholder's custom geometry/masking
                        shape.insert_picture(image_stream)
                        print(f"  ✓ Filled picture placeholder with image from {img_data['name']}")
                        
                    except Exception as e:
                        print(f"  ⚠️  Could not fill picture placeholder: {e}")
                        continue
    
    def _copy_footer_elements(self, template_slide, new_slide, footer_info: Dict[str, str]):
        """Copy footer, slide number, and other non-content placeholders from template"""
        from lxml import etree
        
        # Build the footer text
        footer_text = f"{footer_info['title']} / {footer_info['month']} {footer_info['day']}, {footer_info['year']} / © {footer_info['year']} AT&T Intellectual Property - AT&T Proprietary (Internal Use Only)"
        
        # Find footer and slide number shapes in template
        for shape in template_slide.shapes:
            try:
                # Try to get placeholder_format (will raise exception if not a placeholder)
                ph_type = shape.placeholder_format.type
                ph_type_str = str(ph_type)
                
                # Copy footer and slide number placeholders
                if 'FOOTER' in ph_type_str or 'SLIDE_NUMBER' in ph_type_str:
                    try:
                        # Deep copy the shape element using lxml
                        shape_copy = etree.fromstring(etree.tostring(shape.element))
                        
                        # If this is a footer, update its text
                        if 'FOOTER' in ph_type_str:
                            # Find and update text elements in the XML
                            for text_elem in shape_copy.iter():
                                if text_elem.tag.endswith('}t'):  # Text element
                                    text_elem.text = footer_text
                                    break
                            print(f"  ✓ Copied and updated {shape.name}")
                        else:
                            print(f"  ✓ Copied {shape.name} from template")
                        
                        new_slide.shapes._spTree.append(shape_copy)
                    except Exception as e:
                        print(f"  ⚠️  Could not copy {shape.name}: {e}")
            except ValueError:
                # Not a placeholder, skip
                continue
    
    def _copy_non_placeholder_text_boxes(self, template_slide, new_slide):
        """Copy non-placeholder shapes from template (text boxes, icons, graphics)"""
        from lxml import etree
        
        # Find shapes that are NOT placeholders in template
        for shape in template_slide.shapes:
            # Skip if it's a placeholder
            try:
                _ = shape.placeholder_format
                # If we get here, it's a placeholder - skip it
                continue
            except ValueError:
                # Not a placeholder - this is what we want to copy
                pass
            
            # Copy text boxes (like enumeration numbers)
            if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                text = shape.text.strip()
                # Only copy simple text boxes (like numbers), not complex content
                if text and len(text) < 20:  # Simple text like "1", "2", "3", "4"
                    try:
                        # Deep copy the shape element
                        shape_copy = etree.fromstring(etree.tostring(shape.element))
                        new_slide.shapes._spTree.append(shape_copy)
                        print(f"  ✓ Copied text box: '{text}'")
                    except Exception as e:
                        print(f"  ⚠️  Could not copy text box '{text}': {e}")
            
            # Copy graphics/icons (non-placeholder images) - requires proper image part copying
            elif hasattr(shape, 'image') and 'Graphic' in shape.name:
                try:
                    # Get the image from the shape
                    image_blob = shape.image.blob
                    
                    # Add the image to the new slide with same position/size
                    pic = new_slide.shapes.add_picture(
                        BytesIO(image_blob),
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                    print(f"  ✓ Copied graphic: {shape.name}")
                except Exception as e:
                    print(f"  ⚠️  Could not copy graphic {shape.name}: {e}")
    
    def _copy_placeholder_positions(self, template_slide, new_slide):
        """Copy placeholder positions from template to new slide (layout defaults may differ)"""
        # Build map of template placeholders by index
        template_positions = {}
        for shape in template_slide.shapes:
            try:
                idx = shape.placeholder_format.idx
                template_positions[idx] = {
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
            except ValueError:
                # Not a placeholder
                continue
        
        # Apply positions to matching placeholders in new slide
        for shape in new_slide.shapes:
            try:
                idx = shape.placeholder_format.idx
                if idx in template_positions:
                    pos = template_positions[idx]
                    shape.left = pos['left']
                    shape.top = pos['top']
                    shape.width = pos['width']
                    shape.height = pos['height']
            except ValueError:
                # Not a placeholder
                continue
    
    def inspect_presentation(self, pptx_path: str) -> Dict[str, Any]:
        """Inspect presentation and return detailed info"""
        try:
            prs = Presentation(pptx_path)
            
            slides_data = []
            for idx, slide in enumerate(prs.slides):
                slide_info = {
                    'slide_number': idx + 1,
                    'layout_name': slide.slide_layout.name,
                    'shapes': []
                }
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_info['shapes'].append({
                                'name': shape.name,
                                'text': text
                            })
                
                slides_data.append(slide_info)
            
            return {
                'file': pptx_path,
                'total_slides': len(prs.slides),
                'slides': slides_data
            }
        
        except Exception as e:
            return {
                'file': pptx_path,
                'error': str(e)
            }
    
    def validate_presentation(self, pptx_path: str, output_json: bool = False) -> Dict[str, Any]:
        """
        Validate rendered PPTX against quality rules.
        Returns JSON report with per-slide issues.
        """
        try:
            prs = Presentation(pptx_path)
            
            results = {
                "file": pptx_path,
                "total_slides": len(prs.slides),
                "validation_results": {
                    "overall_status": "pass",
                    "slides": []
                }
            }
            
            for idx, slide in enumerate(prs.slides, start=1):
                issues = []
                
                # Check 1: Missing diagrams (check for actual embedded images)
                layout_name = slide.slide_layout.name
                picture_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                
                # Slides that typically should have diagrams
                # Exclude "Logo End Slide" (12_final_slide) - logo-only, no diagrams expected
                diagram_expected_layouts = ["Blank", "End Slide"]
                is_logo_only_slide = "Logo End Slide" in layout_name
                
                if any(name in layout_name for name in diagram_expected_layouts) and not is_logo_only_slide:
                    if not picture_shapes:
                        issues.append({
                            "type": "missing_embedded_image",
                            "severity": "critical",
                            "message": f"Slide {idx} layout '{layout_name}' has no embedded images",
                            "layout": layout_name,
                            "expected": "At least one embedded image",
                            "actual": "Zero images found"
                        })
                    else:
                        # Verify images are not zero-size (another silent failure mode)
                        for pic in picture_shapes:
                            if pic.width == 0 or pic.height == 0:
                                issues.append({
                                    "type": "zero_size_image",
                                    "severity": "critical",
                                    "message": f"Embedded image '{pic.name}' has zero dimensions",
                                    "image_name": pic.name,
                                    "width": pic.width,
                                    "height": pic.height
                                })
                
                # Check 2: Diagram overlap with footer zone (y > 6.5 inches)
                FOOTER_ZONE_Y = int(FOOTER_ZONE_TOP_INCHES * EMUS_PER_INCH)
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        shape_bottom = shape.top + shape.height
                        if shape_bottom > FOOTER_ZONE_Y:
                            issues.append({
                                "type": "diagram_overlap",
                                "severity": "critical",
                                "message": f"Diagram extends into footer zone (y+height > 6.5\")",
                                "diagram_bottom_inches": round(shape_bottom / EMUS_PER_INCH, 2),
                                "suggested_fix": "Reduce diagram height or adjust y position"
                            })
                
                # Check 2b: Diagram sizing quality (scaling and area coverage)
                # Apply to diagram-only slides (Blank, End Slide) where diagrams should maximize space
                layout_name = slide.slide_layout.name
                DIAGRAM_ONLY_LAYOUTS = ["Blank", "End Slide"]
                
                if any(name in layout_name for name in DIAGRAM_ONLY_LAYOUTS):
                    picture_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                    
                    for shape in picture_shapes:
                        # Convert dimensions from EMUs to inches
                        width_inches = shape.width / EMUS_PER_INCH
                        height_inches = shape.height / EMUS_PER_INCH
                        
                        # Get diagram sizing thresholds from validation config
                        MAX_WIDTH = self.validation.diagram_safe_area_width
                        MAX_HEIGHT = self.validation.diagram_safe_area_height
                        SCALING_THRESHOLD = self.validation.diagram_scaling_threshold
                        COVERAGE_THRESHOLD = self.validation.diagram_coverage_threshold
                        
                        # Calculate metrics
                        scaling_ratio = min(width_inches / MAX_WIDTH, height_inches / MAX_HEIGHT)
                        area_coverage = (width_inches * height_inches) / (MAX_WIDTH * MAX_HEIGHT)
                        
                        # Check 1: Scaling ratio - CRITICAL if below threshold
                        if scaling_ratio < SCALING_THRESHOLD:
                            issues.append({
                                "type": "diagram_undersized",
                                "severity": "critical",
                                "message": f"Diagram scaled to only {scaling_ratio*100:.0f}% of available space (threshold: {SCALING_THRESHOLD*100:.0f}%)",
                                "scaling_ratio": round(scaling_ratio, 3),
                                "diagram_size_inches": f"{width_inches:.2f}×{height_inches:.2f}",
                                "max_size_inches": f"{MAX_WIDTH}×{MAX_HEIGHT}",
                                "suggested_fix": f"Adjust diagram aspect ratio to match slide shape (optimal: {MAX_WIDTH}:{MAX_HEIGHT} = 2.18:1). Consider: 1) change orientation (horizontal→vertical), 2) different diagram type, or 3) restructure content"
                            })
                        
                        # Check 2: Area coverage - CRITICAL if below threshold
                        if area_coverage < COVERAGE_THRESHOLD:
                            issues.append({
                                "type": "diagram_low_coverage",
                                "severity": "critical",
                                "message": f"Diagram uses only {area_coverage*100:.0f}% of available space (threshold: {COVERAGE_THRESHOLD*100:.0f}%)",
                                "area_coverage": round(area_coverage, 3),
                                "diagram_size_inches": f"{width_inches:.2f}×{height_inches:.2f}",
                                "max_size_inches": f"{MAX_WIDTH}×{MAX_HEIGHT}",
                                "suggested_fix": f"Adjust diagram aspect ratio to better match slide (optimal: {MAX_WIDTH}:{MAX_HEIGHT} = 2.18:1)"
                            })
                
                # Check 3: Text overlap detection (exclude divider slides where footer/slide# naturally overlap)
                DIVIDER_LAYOUTS = ["Divider", "Cover w/ Subtitle"]
                layout_name = slide.slide_layout.name
                
                # Skip overlap detection for divider slides (footer/slide# naturally overlap by design)
                if not any(div_name in layout_name for div_name in DIVIDER_LAYOUTS):
                    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
                    
                    # Exclude footer and slide number placeholders from overlap checks
                    IGNORE_PLACEHOLDERS = ["Footer Placeholder", "Slide Number Placeholder"]
                    content_shapes = [
                        s for s in text_shapes 
                        if not any(ignore in s.name for ignore in IGNORE_PLACEHOLDERS)
                    ]
                    
                    for i, shape1 in enumerate(content_shapes):
                        for shape2 in content_shapes[i+1:]:
                            if self._shapes_overlap(shape1, shape2):
                                issues.append({
                                    "type": "text_overlap",
                                    "severity": "critical",
                                    "message": f"Content text shapes '{shape1.name}' and '{shape2.name}' overlap",
                                    "shape1": shape1.name,
                                    "shape2": shape2.name,
                                    "slide_layout": layout_name
                                })
                    
                    # Check 3b: Title overlapping with body content (common on multi-column slides)
                    title_shapes = [s for s in slide.shapes if "Title" in s.name and s.has_text_frame]
                    body_shapes = [
                        s for s in slide.shapes 
                        if s.has_text_frame and "Title" not in s.name 
                        and s.text.strip()
                        and "Footer" not in s.name 
                        and "Slide Number" not in s.name
                    ]
                    
                    for title in title_shapes:
                        title_bottom = title.top + title.height
                        for body in body_shapes:
                            body_top = body.top
                            # Check if title extends below where body content starts
                            if title_bottom > body_top:
                                overlap_inches = (title_bottom - body_top) / EMUS_PER_INCH
                                if overlap_inches > 0.1:  # More than 0.1" overlap
                                    issues.append({
                                        "type": "title_body_overlap",
                                        "severity": "critical",
                                        "message": f"Title '{title.name}' extends into body content region by {overlap_inches:.2f}\"",
                                        "title_shape": title.name,
                                        "body_shape": body.name,
                                        "overlap_inches": round(overlap_inches, 2),
                                        "suggested_fix": "Reduce title font size or use shorter title text"
                                    })
                
                # Check 4: Text brevity (check if text fills > 80% of text frame)
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        text = shape.text.strip()
                        # Rough heuristic: if text is very long, flag it
                        if len(text) > 200 and "statement" in shape.name.lower():
                            issues.append({
                                "type": "text_too_long",
                                "severity": "warning",
                                "message": f"Text in '{shape.name}' appears unnaturally long ({len(text)} chars)",
                                "character_count": len(text),
                                "suggested_max": 200
                            })
                
                # Check 5: Programmatic diagram quality (placeholder - would need shape analysis)
                # This is a simplified check - full implementation would analyze connectors and labels
                connector_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
                if connector_shapes:
                    # Just flag that connectors exist and need manual review for now
                    pass  # Full implementation would check connector endpoints vs shape edges
                
                # Check 6: Detect ellipses truncation (should never happen)
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        text = shape.text.strip()
                        if "..." in text or "…" in text:
                            issues.append({
                                "type": "text_truncated",
                                "severity": "critical",
                                "message": f"Text in '{shape.name}' contains ellipses truncation",
                                "shape_name": shape.name,
                                "text_preview": text[:50] + "..." if len(text) > 50 else text,
                                "suggested_fix": "Rewrite content to fit without truncation"
                            })
                
                if issues:
                    results["validation_results"]["overall_status"] = "fail"
                    results["validation_results"]["slides"].append({
                        "slide_number": idx,
                        "issues": issues
                    })
            
            return results
        
        except Exception as e:
            return {
                "file": pptx_path,
                "error": str(e),
                "validation_results": {
                    "overall_status": "error",
                    "slides": []
                }
            }
    
    def _shapes_overlap(self, shape1, shape2) -> bool:
        """Check if two shapes overlap"""
        # Get bounding boxes
        left1, top1 = shape1.left, shape1.top
        right1, bottom1 = left1 + shape1.width, top1 + shape1.height
        
        left2, top2 = shape2.left, shape2.top
        right2, bottom2 = left2 + shape2.width, top2 + shape2.height
        
        # Check for overlap
        return not (right1 < left2 or right2 < left1 or bottom1 < top2 or bottom2 < top1)


def main():
    import argparse
    
    parser = argparse.ArgumentParser(description='PPT-Maker Engine')
    parser.add_argument('command', choices=['create', 'inspect', 'validate'], help='Command to execute')
    parser.add_argument('path', help='Output path for create, or input path for inspect/validate')
    parser.add_argument('--assets-folder', default='assets', help='Folder for persisting YAML and images (default: assets)')
    parser.add_argument('--json', action='store_true', help='Output validation results as JSON')
    parser.add_argument('--template', 
                       default='.cdo-aifc/templates/09-documentation-requirements/ppt-maker/ppt-maker-template-library.pptx',
                       help='Path to PowerPoint template (default: EAIFC location)')
    
    args = parser.parse_args()
    
    if args.command == "create":
        if sys.stdin.isatty():
            print("Error: No YAML content provided on stdin")
            print("Usage: cat content.yaml | python ppt_engine_focused.py create <output.pptx> [--assets-folder <path>]")
            return
        
        yaml_content = sys.stdin.read()
        
        engine = PPTEngine(args.template, assets_folder=args.assets_folder)
        result_path = engine.create_from_yaml_spec(yaml_content, args.path)
        print(f"\n✅ Presentation created: {result_path}")
    
    elif args.command == "inspect":
        engine = PPTEngine(args.template, assets_folder=args.assets_folder)
        result = engine.inspect_presentation(args.path)
        
        import json
        print(json.dumps(result, indent=2))
    
    elif args.command == "validate":
        engine = PPTEngine(args.template, assets_folder=args.assets_folder)
        result = engine.validate_presentation(args.path, output_json=args.json)
        
        if args.json:
            import json
            print(json.dumps(result, indent=2))
        else:
            # Human-readable output
            print(f"\n📋 Validation Report: {args.path}")
            print(f"Total Slides: {result['total_slides']}")
            print(f"Overall Status: {'✅ PASS' if result['validation_results']['overall_status'] == 'pass' else '❌ FAIL'}\n")
            
            if result['validation_results']['slides']:
                print("Issues Found:")
                for slide in result['validation_results']['slides']:
                    print(f"\n  Slide {slide['slide_number']}:")
                    for issue in slide['issues']:
                        severity_icon = "🔴" if issue['severity'] == 'critical' else "⚠️"
                        print(f"    {severity_icon} [{issue['type']}] {issue['message']}")
            else:
                print("✅ No issues found!")


if __name__ == "__main__":
    main()